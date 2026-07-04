"""
document_generator.py
Generates Word DOCX and PowerPoint PPTX files from campaign artifacts.
Uses python-docx and python-pptx — no LLM required.
"""
import io
import re
import logging
import random
import requests as http_requests
from typing import Optional, Dict, List

logger = logging.getLogger("document_generator")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _fetch_bytes(url: str) -> Optional[bytes]:
    """Download bytes from an HTTPS URL. Returns None on failure."""
    if not url:
        return None
    try:
        resp = http_requests.get(url, timeout=15)
        resp.raise_for_status()
        return resp.content
    except Exception as e:
        logger.warning(f"Failed to fetch bytes from {url}: {e}")
        return None


def _parse_markdown_paragraphs(md: str) -> List[Dict]:
    """
    Very lightweight Markdown → list-of-blocks parser.
    Returns a list of dicts: {type: 'h1'|'h2'|'h3'|'body'|'bullet', text: str}
    """
    blocks = []
    for line in md.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("### "):
            blocks.append({"type": "h3", "text": stripped[4:]})
        elif stripped.startswith("## "):
            blocks.append({"type": "h2", "text": stripped[3:]})
        elif stripped.startswith("# "):
            blocks.append({"type": "h1", "text": stripped[2:]})
        elif stripped.startswith("- ") or stripped.startswith("* "):
            blocks.append({"type": "bullet", "text": stripped[2:]})
        else:
            # Strip bold/italic markers for plain text
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", stripped)
            text = re.sub(r"\*(.+?)\*", r"\1", text)
            blocks.append({"type": "body", "text": text})
    return blocks


# ---------------------------------------------------------------------------
# DOCX Generator
# ---------------------------------------------------------------------------

def generate_docx(
    campaign_name: str,
    company_name: str,
    logo_url: Optional[str],
    blog_post_md: str,
    press_release_md: str,
    longform_md: str,
    blog_hero_url: Optional[str],
    editorial_url: Optional[str],
    slide_background_url: Optional[str],
    content_card_url: Optional[str],
) -> bytes:
    """
    Builds a Word DOCX with:
    - Company logo in header
    - Campaign title
    - All 3 text artifacts as formatted rich text (headings, bullets, body)
    - All 4 images inline with captions
    Returns raw DOCX bytes.
    """
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ── Page margins ────────────────────────────────────────────────────────
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    # ── Header with logo ────────────────────────────────────────────────────
    header = doc.sections[0].header
    header_para = header.paragraphs[0]
    header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    logo_bytes = _fetch_bytes(logo_url) if logo_url else None
    if logo_bytes:
        logo_run = header_para.add_run()
        logo_run.add_picture(io.BytesIO(logo_bytes), width=Inches(1.2))
    else:
        header_para.add_run(company_name).bold = True

    # ── Campaign Title ───────────────────────────────────────────────────────
    title_para = doc.add_heading(campaign_name, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subtitle = doc.add_paragraph(f"{company_name} · Campaign Package")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    doc.add_paragraph()  # spacer

    # ── Helper: add a labelled image ────────────────────────────────────────
    def add_image_section(label: str, img_url: Optional[str], width_inches: float = 5.5):
        img_bytes = _fetch_bytes(img_url)
        if img_bytes:
            doc.add_heading(label, level=2)
            img_para = doc.add_paragraph()
            img_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            img_run = img_para.add_run()
            img_run.add_picture(io.BytesIO(img_bytes), width=Inches(width_inches))
            doc.add_paragraph()

    # ── Section 1: Images ────────────────────────────────────────────────────
    doc.add_heading("Visual Assets", level=1)
    add_image_section("Blog / PR Hero Image", blog_hero_url, 5.5)
    add_image_section("Editorial Image", editorial_url, 4.5)
    add_image_section("Slide Background", slide_background_url, 5.5)
    add_image_section("Content Card", content_card_url, 3.0)

    doc.add_page_break()

    # ── Helper: render markdown section ─────────────────────────────────────
    def add_markdown_section(section_title: str, md_text: str):
        if not md_text or not md_text.strip():
            return
        doc.add_heading(section_title, level=1)
        blocks = _parse_markdown_paragraphs(md_text)
        for block in blocks:
            if block["type"] == "h1":
                doc.add_heading(block["text"], level=2)
            elif block["type"] == "h2":
                doc.add_heading(block["text"], level=3)
            elif block["type"] == "h3":
                doc.add_heading(block["text"], level=4)
            elif block["type"] == "bullet":
                p = doc.add_paragraph(block["text"], style="List Bullet")
            else:
                p = doc.add_paragraph(block["text"])
                p.style.font.size = Pt(10)
        doc.add_paragraph()

    # ── Section 2: Text Artifacts ────────────────────────────────────────────
    add_markdown_section("Blog Post", blog_post_md)
    doc.add_page_break()
    add_markdown_section("Press Release", press_release_md)
    doc.add_page_break()
    add_markdown_section("Long-Form Editorial", longform_md)

    # ── Serialize to bytes ───────────────────────────────────────────────────
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# PPTX Generator
# ---------------------------------------------------------------------------

# 10 slide layout configurations: (layout_fn_name, accent_color_hex)
# Each layout is a function that populates one slide programmatically.
_ACCENT_COLORS = [
    (0x0F, 0x17, 0x2A),  # slate-950
    (0x4F, 0x46, 0xE5),  # indigo-600
    (0x0E, 0xA5, 0xE9),  # sky-500
    (0x10, 0xB9, 0x81),  # emerald-500
    (0x1E, 0x40, 0xAF),  # blue-800
]


def _pptx_add_logo(slide, logo_bytes: Optional[bytes], slide_width, slide_height):
    """Place logo in top-right corner if available."""
    from pptx.util import Inches, Pt
    if not logo_bytes:
        return
    try:
        left  = slide_width - Inches(1.4)
        top   = Inches(0.1)
        slide.shapes.add_picture(io.BytesIO(logo_bytes), left, top, width=Inches(1.2))
    except Exception as e:
        logger.warning(f"Failed to add logo to slide: {e}")


def _pptx_set_bg(slide, bg_bytes: Optional[bytes]):
    """Set slide background to a full-bleed image."""
    from pptx.util import Inches
    if not bg_bytes:
        return
    try:
        slide.background.fill.user_picture(io.BytesIO(bg_bytes))
    except Exception as e:
        logger.warning(f"Failed to set slide background natively: {e}")
        try:
            # Fallback: add picture shape and move to back of the shape tree
            pic = slide.shapes.add_picture(
                io.BytesIO(bg_bytes), 0, 0, width=Inches(13.33), height=Inches(7.5)
            )
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)
        except Exception as e2:
            logger.warning(f"Fallback background set also failed: {e2}")


def generate_pptx(
    campaign_name: str,
    company_name: str,
    logo_url: Optional[str],
    slide_bg_url: Optional[str],
    blog_hero_url: Optional[str],
    editorial_url: Optional[str],
    content_card_url: Optional[str],
    blog_post_md: str,
    press_release_md: str,
    master_pillars: Optional[List[Dict]] = None,
) -> bytes:
    """
    Builds a PowerPoint PPTX with 10 slides of varied layouts.
    Slide background on most slides; company logo in corner.
    Returns raw PPTX bytes.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    SH = prs.slide_height

    blank_layout = prs.slide_layouts[6]  # completely blank

    # Pre-fetch assets once
    logo_bytes   = _fetch_bytes(logo_url)
    bg_bytes     = _fetch_bytes(slide_bg_url)
    hero_bytes   = _fetch_bytes(blog_hero_url)
    ed_bytes     = _fetch_bytes(editorial_url)
    card_bytes   = _fetch_bytes(content_card_url)

    # Extract text snippets from markdown
    blog_blocks   = _parse_markdown_paragraphs(blog_post_md)
    pr_blocks     = _parse_markdown_paragraphs(press_release_md)
    blog_bullets  = [b["text"] for b in blog_blocks if b["type"] == "bullet"][:5]
    pr_bullets    = [b["text"] for b in pr_blocks  if b["type"] == "bullet"][:5]
    blog_body_snippets = [b["text"] for b in blog_blocks if b["type"] == "body"][:3]
    pr_body_snippets   = [b["text"] for b in pr_blocks  if b["type"] == "body"][:2]

    pillars = master_pillars or [
        {"name": "Driver Exoneration & Insurance Reduction"},
        {"name": "Proactive In-Cab AI Coaching"},
        {"name": "Asset Optimization & Fuel Savings"},
    ]

    accent = _ACCENT_COLORS[1]  # indigo

    # ── Utility: add text box ────────────────────────────────────────────────
    def add_text_box(slide, text: str, left, top, width, height,
                     font_size=24, bold=False, color=(0xFF, 0xFF, 0xFF),
                     align=PP_ALIGN.LEFT, wrap=True):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = PptxRGB(*color)
        return txBox

    def add_rect(slide, left, top, width, height, fill_rgb=(0x0F, 0x17, 0x2A), alpha=None):
        from pptx.oxml.ns import qn as pqn
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = PptxRGB(*fill_rgb)
        shape.line.fill.background()
        return shape

    # ── SLIDE 1: Title Slide ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    # Slate/indigo text styling designed for light backgrounds
    add_text_box(s, company_name,     Inches(1), Inches(2.0), Inches(11), Inches(1),   font_size=20, bold=True,  color=(0x4F, 0x46, 0xE5), align=PP_ALIGN.CENTER)
    add_text_box(s, campaign_name,    Inches(1), Inches(2.9), Inches(11), Inches(1.5), font_size=40, bold=True,  color=(0x0F, 0x17, 0x2A), align=PP_ALIGN.CENTER)
    add_text_box(s, "Campaign Package", Inches(1), Inches(4.2), Inches(11), Inches(0.8), font_size=18, color=(0x47, 0x55, 0x69), align=PP_ALIGN.CENTER)
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 2: Blog Hero Full-Bleed ────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if hero_bytes:
        pic = s.shapes.add_picture(io.BytesIO(hero_bytes), 0, 0, width=SW, height=SH)
        from lxml import etree
        sp = s.shapes._spTree
        sp.remove(pic._element); sp.insert(2, pic._element)
    add_rect(s, 0, Inches(5.8), SW, Inches(1.7), (0x0A, 0x0F, 0x1E))
    add_text_box(s, "Campaign Visual Overview", Inches(0.5), Inches(6.0), Inches(12), Inches(0.7), font_size=22, bold=True, color=(0xFF, 0xFF, 0xFF))
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 3: Pillar #1 ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    add_rect(s, 0, 0, Inches(5.5), SH, (0x0F, 0x17, 0x2A))
    add_text_box(s, "Core Value Proposition", Inches(0.3), Inches(1.0), Inches(4.8), Inches(0.8), font_size=13, color=(0x94, 0xA3, 0xB8))
    p1 = pillars[0]["name"] if pillars else "Driver Safety"
    add_text_box(s, p1, Inches(0.3), Inches(1.7), Inches(4.8), Inches(1.5), font_size=26, bold=True, color=(0xFF, 0xFF, 0xFF))
    body = blog_body_snippets[0] if blog_body_snippets else f"{company_name} helps fleet operators reduce risk and protect drivers with AI-powered video telematics."
    add_text_box(s, body, Inches(0.3), Inches(3.2), Inches(4.8), Inches(2.5), font_size=14, color=(0xCB, 0xD5, 0xE1))
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 4: Pillar #2 with editorial image ──────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    add_rect(s, Inches(6.8), 0, Inches(6.5), SH, (0x0F, 0x17, 0x2A))
    if ed_bytes:
        s.shapes.add_picture(io.BytesIO(ed_bytes), 0, 0, width=Inches(6.6), height=SH)
    p2 = pillars[1]["name"] if len(pillars) > 1 else "AI Coaching"
    add_text_box(s, "Value Driver", Inches(7.0), Inches(1.0), Inches(6.0), Inches(0.7), font_size=13, color=(0x94, 0xA3, 0xB8))
    add_text_box(s, p2, Inches(7.0), Inches(1.6), Inches(6.0), Inches(1.3), font_size=26, bold=True, color=(0xFF, 0xFF, 0xFF))
    body2 = blog_body_snippets[1] if len(blog_body_snippets) > 1 else "Real-time coaching alerts reduce harsh events and improve driver behaviour on every route."
    add_text_box(s, body2, Inches(7.0), Inches(3.0), Inches(6.0), Inches(2.5), font_size=14, color=(0xCB, 0xD5, 0xE1))
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 5: Blog Post Key Points ────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    add_rect(s, 0, 0, SW, Inches(1.5), accent)
    add_text_box(s, "Blog Post — Key Insights", Inches(0.5), Inches(0.35), Inches(12), Inches(0.9), font_size=24, bold=True, color=(0xFF, 0xFF, 0xFF))
    bullets = blog_bullets or ["Fleet telematics improves driver safety", "Real-time video exonerates drivers", "Operational ROI through data intelligence"]
    for i, bullet in enumerate(bullets[:5]):
        add_text_box(s, f"• {bullet}", Inches(0.7), Inches(1.7 + i * 0.95), Inches(11.5), Inches(0.9), font_size=16, color=(0x1E, 0x29, 0x3B))
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 6: Pillar #3 ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    add_rect(s, 0, 0, Inches(5.5), SH, (0x1E, 0x40, 0xAF))
    p3 = pillars[2]["name"] if len(pillars) > 2 else "Operational Efficiency"
    add_text_box(s, "Operational Focus", Inches(0.3), Inches(1.0), Inches(4.8), Inches(0.8), font_size=13, color=(0xBF, 0xDB, 0xFE))
    add_text_box(s, p3, Inches(0.3), Inches(1.7), Inches(4.8), Inches(1.5), font_size=26, bold=True, color=(0xFF, 0xFF, 0xFF))
    body3 = blog_body_snippets[2] if len(blog_body_snippets) > 2 else "Optimize fuel usage, reduce idle time, and cut operational costs across your entire fleet."
    add_text_box(s, body3, Inches(0.3), Inches(3.2), Inches(4.8), Inches(2.5), font_size=14, color=(0xBF, 0xDB, 0xFE))
    if card_bytes:
        s.shapes.add_picture(io.BytesIO(card_bytes), Inches(5.8), Inches(0.5), width=Inches(7.0), height=Inches(6.5))
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 7: Press Release Summary ──────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    add_rect(s, 0, 0, SW, Inches(1.5), (0x0F, 0x17, 0x2A))
    add_text_box(s, "Press Release Summary", Inches(0.5), Inches(0.35), Inches(12), Inches(0.9), font_size=24, bold=True, color=(0xFF, 0xFF, 0xFF))
    pr_points = pr_bullets or pr_body_snippets or ["Official announcement to media and industry analysts.", "Verified data points support the release claims."]
    for i, pt in enumerate(pr_points[:5]):
        add_text_box(s, f"• {pt}", Inches(0.7), Inches(1.7 + i * 0.95), Inches(11.5), Inches(0.9), font_size=16, color=(0x1E, 0x29, 0x3B))
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 8: Content Card Focus ──────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    add_rect(s, 0, 0, SW, Inches(1.5), (0x0F, 0x17, 0x2A))
    add_text_box(s, "Product Focus", Inches(0.5), Inches(0.35), Inches(8), Inches(0.9), font_size=24, bold=True, color=(0xFF, 0xFF, 0xFF))
    if card_bytes:
        s.shapes.add_picture(io.BytesIO(card_bytes), Inches(0.4), Inches(1.7), width=Inches(4.5), height=Inches(4.5))
    add_text_box(s, "Fleet Dashcam Solution", Inches(5.3), Inches(1.9), Inches(7.3), Inches(0.9), font_size=22, bold=True, color=(0x1E, 0x29, 0x3B))
    features = ["Dual-facing HD cameras", "Real-time driver coaching alerts", "Cloud-connected event upload", "Insurance & compliance reporting"]
    for i, feat in enumerate(features):
        add_text_box(s, f"✓ {feat}", Inches(5.3), Inches(2.9 + i * 0.8), Inches(7.3), Inches(0.75), font_size=15, color=(0x1E, 0x29, 0x3B))
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 9: CTA Slide ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    add_text_box(s, "Ready to Transform Your Fleet?", Inches(1), Inches(2.0), Inches(11), Inches(1.3), font_size=34, bold=True, color=(0x0F, 0x17, 0x2A), align=PP_ALIGN.CENTER)
    add_text_box(s, "Request a Free Hardware Pilot or Book a Live Demo", Inches(1), Inches(3.4), Inches(11), Inches(0.9), font_size=20, color=(0x47, 0x55, 0x69), align=PP_ALIGN.CENTER)
    add_rect(s, Inches(4.5), Inches(4.6), Inches(4.3), Inches(0.9), (0x4F, 0x46, 0xE5))
    add_text_box(s, "Book a Live Demo →", Inches(4.5), Inches(4.65), Inches(4.3), Inches(0.8), font_size=18, bold=True, color=(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
    _pptx_add_logo(s, logo_bytes, SW, SH)

    # ── SLIDE 10: Thank You ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    if bg_bytes:
        _pptx_set_bg(s, bg_bytes)
    if logo_bytes:
        s.shapes.add_picture(io.BytesIO(logo_bytes), Inches(5.5), Inches(1.5), width=Inches(2.3))
    add_text_box(s, "Thank You", Inches(1), Inches(3.5), Inches(11), Inches(1.2), font_size=42, bold=True, color=(0x0F, 0x17, 0x2A), align=PP_ALIGN.CENTER)
    add_text_box(s, company_name, Inches(1), Inches(4.8), Inches(11), Inches(0.7), font_size=18, color=(0x4F, 0x46, 0xE5), align=PP_ALIGN.CENTER)
    _pptx_add_logo(s, logo_bytes, SW, SH)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
